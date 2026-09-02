"""
Document Ingestor
─────────────────
Orchestrates the full ingestion pipeline:

    Upload → Parse → OCR (if needed) → PII scrub → Chunk → Embed → Vectorstore

Supports: PDF, DOCX, PPTX, TXT, Markdown
"""
from __future__ import annotations

import hashlib
import logging
import mimetypes
import shutil
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from langchain_core.documents import Document

from config import settings
from ingestion.chunker import DocumentChunker
from ingestion.pii_scrubber import scrub_pii

logger = logging.getLogger(__name__)


# ── Document Parsers ──────────────────────────────────────────────────────────

def _parse_pdf(path: Path) -> list[Document]:
    """Extract text page-by-page using PyMuPDF; fall back to OCR if needed."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError("PyMuPDF not installed. Run: pip install pymupdf")

    docs: list[Document] = []
    pdf = fitz.open(str(path))
    for page_num, page in enumerate(pdf, start=1):
        text = page.get_text("text").strip()

        # Page has no selectable text → OCR
        if not text or len(text) < 20:
            text = _ocr_page(page)

        if text:
            docs.append(Document(
                page_content=text,
                metadata={"page": page_num, "total_pages": len(pdf)},
            ))
    pdf.close()
    return docs


def _ocr_page(page: Any) -> str:  # page: fitz.Page
    """Rasterise a PDF page and run Tesseract OCR."""
    try:
        import pytesseract
        from PIL import Image
        import io

        mat = page.get_pixmap(dpi=200)
        img = Image.open(io.BytesIO(mat.tobytes("png")))
        return pytesseract.image_to_string(img, lang="eng").strip()
    except Exception as exc:
        logger.warning("OCR failed on page %s: %s", getattr(page, "number", "?"), exc)
        return ""


def _parse_docx(path: Path) -> list[Document]:
    try:
        from docx import Document as DocxDocument
    except ImportError:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")

    d = DocxDocument(str(path))
    full_text = "\n\n".join(p.text for p in d.paragraphs if p.text.strip())
    return [Document(page_content=full_text, metadata={"page": 1, "total_pages": 1})]


def _parse_pptx(path: Path) -> list[Document]:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation(str(path))
    docs: list[Document] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={"page": slide_num, "total_pages": len(prs.slides)},
            ))
    return docs


def _parse_text(path: Path) -> list[Document]:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    return [Document(page_content=text, metadata={"page": 1, "total_pages": 1})]


_PARSERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".txt": _parse_text,
    ".md": _parse_text,
}


# ── Ingestor ──────────────────────────────────────────────────────────────────

class DocumentIngestor:
    """
    High-level orchestrator for the document ingestion pipeline.

    Usage::

        from vectorstore.store import VectorStoreManager
        from ingestion import DocumentIngestor

        store = VectorStoreManager()
        ingestor = DocumentIngestor(store)
        result = ingestor.ingest_file(
            file_path=Path("policy.pdf"),
            doc_type="policy",
            access_groups=["hr", "all_staff"],
            uploaded_by="alice@acme.com",
        )
    """

    def __init__(self, vector_store: Any) -> None:
        self._store = vector_store
        self._chunker = DocumentChunker()

    # ─────────────────────────────────────────────────────────────────────────

    def ingest_file(
        self,
        file_path: Path,
        doc_type: str = "general",
        access_groups: list[str] | None = None,
        uploaded_by: str = "system",
        extra_metadata: dict | None = None,
    ) -> dict:
        """
        Full ingestion pipeline for a single file.

        Returns a summary dict with ingestion stats.
        """
        if access_groups is None:
            access_groups = ["all_staff"]

        suffix = file_path.suffix.lower()
        if suffix not in _PARSERS:
            raise ValueError(f"Unsupported file type: {suffix}")

        doc_id = str(uuid.uuid4())
        file_hash = _file_hash(file_path)

        # ── 1. Save to upload directory ──────────────────────────────────────
        dest = settings.UPLOAD_DIR / f"{doc_id}{suffix}"
        settings.UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
        shutil.copy2(file_path, dest)

        # ── 2. Parse ─────────────────────────────────────────────────────────
        logger.info("Parsing %s (%s)", file_path.name, suffix)
        raw_docs = _PARSERS[suffix](dest)

        # ── 3. PII scrub ─────────────────────────────────────────────────────
        for d in raw_docs:
            d.page_content = scrub_pii(d.page_content)

        # ── 4. Attach base metadata ───────────────────────────────────────────
        base_meta = {
            "doc_id": doc_id,
            "source_filename": file_path.name,
            "doc_type": doc_type,
            "access_groups": ",".join(access_groups),  # stored as CSV string
            "uploaded_by": uploaded_by,
            "ingestion_timestamp": datetime.now(timezone.utc).isoformat(),
            "file_hash": file_hash,
        }
        if extra_metadata:
            base_meta.update(extra_metadata)

        for d in raw_docs:
            d.metadata.update(base_meta)

        # ── 5. Chunk ──────────────────────────────────────────────────────────
        chunks = self._chunker.chunk_documents(raw_docs)
        logger.info("Chunked %s → %d chunks", file_path.name, len(chunks))

        # ── 6. Add to vector store ────────────────────────────────────────────
        self._store.add_documents(chunks)
        logger.info("Indexed %d chunks for doc_id=%s", len(chunks), doc_id)

        return {
            "doc_id": doc_id,
            "filename": file_path.name,
            "doc_type": doc_type,
            "access_groups": access_groups,
            "pages_parsed": len(raw_docs),
            "chunks_indexed": len(chunks),
            "file_hash": file_hash,
            "ingested_at": base_meta["ingestion_timestamp"],
            "uploaded_by": uploaded_by,
        }

    def delete_document(self, doc_id: str) -> int:
        """Remove all chunks for a given document from the vector store."""
        return self._store.delete_by_doc_id(doc_id)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _file_hash(path: Path) -> str:
    """SHA-256 hash for deduplication."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()
